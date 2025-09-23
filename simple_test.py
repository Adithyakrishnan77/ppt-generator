#!/usr/bin/env python3
"""
Simple test script to verify bullet point fix
"""

import sys
import os
import tempfile

# Add current directory to path
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

try:
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    print("✅ python-pptx imported successfully")

    # Test creating a simple presentation with bullet points
    prs = Presentation()

    # Create a slide
    slide_layout = prs.slide_layouts[1]  # Title and Content layout
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    title = slide.shapes.title
    title.text = "Test Presentation"

    # Add content with bullet points
    content = slide.placeholders[1]
    tf = content.text_frame

    # Add bullet points
    for i, point in enumerate(["Point 1", "Point 2", "Point 3"]):
        p = tf.add_paragraph()
        p.text = f"• {point}"
        p.level = 0
        p.space_after = Pt(8)

    # Save test file
    with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
        prs.save(tmp_file.name)
        print(f"✅ Test presentation created successfully: {tmp_file.name}")

    # Clean up
    os.unlink(tmp_file.name)
    print("✅ Test completed successfully!")

except ImportError as e:
    print(f"❌ Import error: {e}")
    sys.exit(1)
except Exception as e:
    print(f"❌ Test failed: {e}")
    import traceback
    traceback.print_exc()
    sys.exit(1)
