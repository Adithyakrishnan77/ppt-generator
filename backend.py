import subprocess
import json
import re
import ast
import os
import threading
import time
from typing import Any, Dict, List, Optional, Tuple
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from functools import lru_cache
import tempfile
import logging

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

class PPTGenerationError(Exception):
    """Custom exception for PPT generation errors"""
    pass

class GenerationStopped(Exception):
    """Exception raised when generation is stopped by user"""
    pass

# --------------------------
# Configuration Management
# --------------------------

@lru_cache(maxsize=1)
def load_config(config_file: str = "prompt_config.json") -> Dict[str, Any]:
    """Load configuration from external JSON file (cached for performance)."""
    
    default_config = {
        "system_prompt": (
            "You are an intelligent slide generator AI that adapts presentation length to content complexity.\n"
            "You must produce ONLY a JSON array inside <json>...</json> tags.\n"
            "Each item MUST be an object: {\"title\": string, \"points\": [string, ...]}.\n"
            "Analyze the user's request and determine the optimal number of slides needed to cover the topic comprehensively.\n\n"
            "Additional Requirements:\n"
            "- Support mathematical content, including LaTeX-style equations (use $...$ inline or $$...$$ for block display).\n"
            "- Include example problems and step-by-step solutions where relevant.\n"
            "- When coding or algorithms are needed, include short code snippets in markdown format (```python ... ```).\n"
            "- Mix concepts, theory, and problem-solving for a comprehensive explanation."
        ),
        "adaptive_prompt_template": (
            "{system}\n\n"
            "User request:\n"
            "{user_request}\n\n"
            "Rules:\n"
            "- Return ONLY a JSON array inside <json>...</json>.\n"
            "- Do not include any text outside the JSON tags.\n"
            "- Slide 1: Title only (no points).\n"
            "- Slides 2+: {min_points}-{max_points} bullet points (50–60 words each).\n"
            "- Use • at the start of every point.\n"
            "- Use LaTeX for equations, e.g., $E = mc^2$, or $$\\int_0^1 x^2 dx = 1/3$$.\n"
            "- Provide example problems and step-by-step solutions if topic is mathematical.\n"
            "- When coding is useful, include code blocks inside points using markdown triple backticks.\n"
            "- Decide the number of slides (3–20) based on topic complexity.\n\n"
        ),
        "min_points": 4,
        "max_points": 12,
        "max_slides": 20,
        "min_slides": 3,
        "model": "qwen3",
        "ollama_timeout": 120  # 2 minutes timeout
    }

    try:
        if os.path.exists(config_file):
            with open(config_file, 'r', encoding='utf-8') as f:
                config = json.load(f)
            # Merge defaults for missing keys
            for key, value in default_config.items():
                if key not in config:
                    config[key] = value
            return config
        else:
            with open(config_file, 'w', encoding='utf-8') as f:
                json.dump(default_config, f, indent=2, ensure_ascii=False)
            logger.info(f"Created default config file: {config_file}")
            return default_config
    except Exception as e:
        logger.warning(f"Error loading config, using defaults: {e}")
        return default_config

# --------------------------
# Content Analysis
# --------------------------

def estimate_slide_count_from_keywords(user_input: str) -> int:
    """Optimized heuristic to estimate slide count based on input complexity."""
    text = user_input.lower()
    
    # Pre-compiled patterns for better performance
    complexity_indicators = {
        'broad': ['overview', 'introduction', 'complete', 'comprehensive', 'full', 'entire'],
        'process': ['process', 'steps', 'workflow', 'procedure', 'implementation'],
        'training': ['training', 'course', 'learning', 'education', 'tutorial']
    }
    
    score = 0.0
    score += sum(2 for word in complexity_indicators['broad'] if word in text)
    score += sum(1.5 for word in complexity_indicators['process'] if word in text)
    score += sum(2 for word in complexity_indicators['training'] if word in text)
    score += text.count(",") * 0.5
    score += len(user_input.split()) // 25  # Word count factor
    
    return max(3, min(18, int(score + 5)))

# --------------------------
# LLM Integration with Stop Support
# --------------------------

class OllamaClient:
    """Thread-safe Ollama client with cancellation support"""
    
    def __init__(self, timeout: int = 120):
        self.timeout = timeout
        self._process = None
        self._stop_event = threading.Event()
    
    def ask_ollama(self, prompt: str, model: str = "qwen3", stop_event: threading.Event = None) -> str:
        """Send prompt to Ollama with cancellation support"""
        if stop_event and stop_event.is_set():
            raise GenerationStopped("Generation stopped by user")
            
        try:
            self._process = subprocess.Popen(
                ["ollama", "run", model],
                stdin=subprocess.PIPE,
                stdout=subprocess.PIPE,
                stderr=subprocess.PIPE,
                text=True,
                encoding="utf-8",
                errors="replace"
            )
            
            # Monitor for stop events during generation
            def check_stop():
                while self._process and self._process.poll() is None:
                    if stop_event and stop_event.is_set():
                        try:
                            self._process.terminate()
                            time.sleep(0.5)
                            if self._process.poll() is None:
                                self._process.kill()
                        except:
                            pass
                        return
                    time.sleep(0.1)
            
            if stop_event:
                stop_thread = threading.Thread(target=check_stop, daemon=True)
                stop_thread.start()
            
            stdout, stderr = self._process.communicate(input=prompt, timeout=self.timeout)
            
            # Check if stopped during communication
            if stop_event and stop_event.is_set():
                raise GenerationStopped("Generation stopped by user")
            
            if self._process.returncode != 0:
                # Don't treat termination due to stop as an error
                if stop_event and stop_event.is_set():
                    raise GenerationStopped("Generation stopped by user")
                raise PPTGenerationError(f"Ollama failed: {stderr or 'Unknown error'}")
                
            return stdout.strip()
            
        except subprocess.TimeoutExpired:
            if self._process:
                self._process.kill()
            raise PPTGenerationError(f"Ollama request timed out after {self.timeout} seconds")
        except FileNotFoundError:
            raise PPTGenerationError("Ollama executable not found. Please install Ollama.")
        except GenerationStopped:
            raise
        except Exception as e:
            if stop_event and stop_event.is_set():
                raise GenerationStopped("Generation stopped by user")
            raise PPTGenerationError(f"Ollama request failed: {str(e)}")
        finally:
            self._process = None

# --------------------------
# JSON Processing
# --------------------------

def extract_and_parse_json(text: str) -> Optional[List[Dict[str, Any]]]:
    """Optimized JSON extraction and parsing with multiple fallback strategies"""
    if not text:
        return None
    
    # Strategy 1: Extract from <json> tags
    json_match = re.search(r"<json>\s*(.*?)\s*</json>", text, re.DOTALL | re.IGNORECASE)
    if json_match:
        json_str = json_match.group(1).strip()
    else:
        # Strategy 2: Extract first JSON array
        start = text.find("[")
        if start == -1:
            return None
        
        depth = 0
        for i, ch in enumerate(text[start:], start=start):
            if ch == "[":
                depth += 1
            elif ch == "]":
                depth -= 1
                if depth == 0:
                    json_str = text[start:i+1].strip()
                    break
        else:
            return None
    
    # Clean common issues
    json_str = (json_str
                .replace(""", '"').replace(""", '"')
                .replace("'", "'").replace("'", "'")
                .replace("\u2013", "-").replace("\u2014", "-"))
    json_str = re.sub(r",\s*([\]\}])", r"\1", json_str)  # Remove trailing commas
    
    # Parse attempts
    for parser in [json.loads, ast.literal_eval]:
        try:
            result = parser(json_str)
            if isinstance(result, list):
                return result
        except:
            continue
    
    return None

def validate_and_coerce_slides(data: List[Any], max_slides: int) -> List[Dict[str, Any]]:
    """Convert parsed data to valid slide format"""
    slides = []
    
    for item in data[:max_slides]:
        if not isinstance(item, dict):
            continue
            
        title = str(item.get("title", "Untitled Slide"))
        points = item.get("points", [])
        
        if not isinstance(points, list):
            points = []
            
        # Clean and validate points
        clean_points = []
        for point in points:
            if isinstance(point, str) and point.strip():
                clean_points.append(point.strip())
        
        slides.append({
            "title": title,
            "points": clean_points
        })
    
    return slides

# --------------------------
# PowerPoint Creation - FIXED
# --------------------------

def create_presentation(
    slide_data: List[Dict[str, Any]],
    template_path: Optional[str] = None,
    font_name: str = "Calibri",
    font_size: int = 12,
    topic : str = "Presentation"
) -> str:
    """Create PowerPoint presentation with FIXED bullet point formatting"""
    
    try:
        # Load template or create new presentation
        if template_path and os.path.exists(template_path):
            prs = Presentation(template_path)
            logger.info(f"Using template: {template_path}")
        else:
            prs = Presentation()
            if template_path:
                logger.warning(f"Template not found: {template_path}, using default")
        
        for idx, slide in enumerate(slide_data):
            # Choose layout - title slide for first, content for others
            if idx == 0:
                layout = prs.slide_layouts[0]  # Title slide layout
            else:
                # Find a content layout (usually index 1, but fallback to 0)
                layout_idx = 1 if len(prs.slide_layouts) > 1 else 0
                layout = prs.slide_layouts[layout_idx]
            
            slide_obj = prs.slides.add_slide(layout)
            
            # Set title
            title_text = slide.get("title", "Untitled Slide")
            if slide_obj.shapes.title:
                title_shape = slide_obj.shapes.title
                title_shape.text = title_text
                
                # Format title
                if title_shape.has_text_frame:
                    for paragraph in title_shape.text_frame.paragraphs:
                        paragraph.alignment = PP_ALIGN.CENTER
                        for run in paragraph.runs:
                            run.font.bold = True
                            run.font.name = font_name
                            run.font.size = Pt(font_size + 8)
                            run.font.color.rgb = RGBColor(31, 56, 148)  # Dark blue
            
            # Add content for non-title slides - FIXED LOGIC
            if idx > 0 and slide.get("points"):
                content_added = False
                
                # Strategy 1: Look for content placeholder by placeholder type
                for shape in slide_obj.shapes:
                    if hasattr(shape, 'has_text_frame') and shape.has_text_frame:
                        # Skip title shapes
                        if shape == slide_obj.shapes.title:
                            continue
                        
                        # Check if this is a content placeholder
                        if hasattr(shape, 'placeholder_format'):
                            try:
                                # This is likely a content placeholder
                                text_frame = shape.text_frame
                                text_frame.clear()  # Clear existing content
                                
                                # RELIABLE BULLET FIX: Manual bullets for all strategies
                                for point_idx, point in enumerate(slide["points"]):
                                    point_text = point.strip()
                                    # Clean any existing bullet symbols
                                    if point_text.startswith('•'):
                                        point_text = point_text[1:].strip()
                                    if point_text.startswith('-'):
                                        point_text = point_text[1:].strip()
                                    
                                    if point_idx == 0:
                                        p = text_frame.paragraphs[0]
                                    else:
                                        p = text_frame.add_paragraph()
                                    
                                    # Add manual bullet - most reliable approach
                                    p.text = f"• {point_text}"
                                    p.level = 0
                                    p.space_after = Pt(6)
                                    
                                    # Format the text runs
                                    for run in p.runs:
                                        run.font.name = font_name
                                        run.font.size = Pt(font_size)
                                        run.font.color.rgb = RGBColor(0, 0, 0)
                                
                                content_added = True
                                break
                            except:
                                continue
                
                # Strategy 2: Fallback - find any text frame that's not the title
                if not content_added:
                    for shape in slide_obj.shapes:
                        if (hasattr(shape, 'has_text_frame') and shape.has_text_frame and 
                            shape != slide_obj.shapes.title):
                            
                            text_frame = shape.text_frame
                            text_frame.clear()
                            
                            # RELIABLE BULLET FIX for Strategy 2
                            for point_idx, point in enumerate(slide["points"]):
                                point_text = point.strip()
                                # Clean any existing bullet symbols
                                if point_text.startswith('•'):
                                    point_text = point_text[1:].strip()
                                if point_text.startswith('-'):
                                    point_text = point_text[1:].strip()
                                
                                if point_idx == 0:
                                    p = text_frame.paragraphs[0]
                                else:
                                    p = text_frame.add_paragraph()
                                
                                # Add manual bullet - consistent with Strategy 1
                                p.text = f"• {point_text}"
                                p.level = 0
                                p.space_after = Pt(6)
                                
                                for run in p.runs:
                                    run.font.name = font_name
                                    run.font.size = Pt(font_size)
                                    run.font.color.rgb = RGBColor(0, 0, 0)
                            
                            content_added = True
                            break
                
                # Strategy 3: Ultimate fallback - add text box if no content area found
                if not content_added:
                    logger.warning(f"No content placeholder found for slide {idx + 1}, adding text box")
                    
                    # Add a text box
                    from pptx.util import Inches
                    left = Inches(1)
                    top = Inches(2)
                    width = Inches(8)
                    height = Inches(5)
                    
                    text_box = slide_obj.shapes.add_textbox(left, top, width, height)
                    text_frame = text_box.text_frame
                    text_frame.clear()
                    
                    for point_idx, point in enumerate(slide["points"]):
                        point_text = point.strip()
                        # Clean any existing bullet symbols
                        if point_text.startswith('•'):
                            point_text = point_text[1:].strip()
                        if point_text.startswith('-'):
                            point_text = point_text[1:].strip()
                        
                        if point_idx == 0:
                            p = text_frame.paragraphs[0]
                        else:
                            p = text_frame.add_paragraph()
                        
                        p.text = f"• {point_text}"
                        p.level = 0
                        p.space_after = Pt(6)
                        
                        for run in p.runs:
                            run.font.name = font_name
                            run.font.size = Pt(font_size)
                            run.font.color.rgb = RGBColor(0, 0, 0)
        
        # Save to temporary file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp_file:
            prs.save(tmp_file.name)
            logger.info(f"Presentation saved to: {tmp_file.name}")
            return tmp_file.name
        # safe_topic = re.sub(r'[^a-zA-Z0-9_-]', '_', topic.strip()) or "Presentation"
        # tmp_dir = tempfile.gettempdir()
        # file_path = os.path.join(tmp_dir, f"{safe_topic}.pptx")

        # prs.save(file_path)
        # logger.info(f"Presentation saved to: {file_path}")
        # return file_path
   
    except Exception as e:
        raise PPTGenerationError(f"Failed to create presentation: {str(e)}")

# --------------------------
# Main Generation Function
# --------------------------

def generate_presentation(
    topic: str,
    template_path: Optional[str] = None,
    font_name: str = "Calibri",
    font_size: int = 12,
    stop_event: Optional[threading.Event] = None
) -> Tuple[str, Dict[str, Any]]:
    """
    Generate PowerPoint presentation from topic with stop support.
    
    Returns:
        Tuple of (file_path, metadata)
    """
    
    if not topic.strip():
        raise PPTGenerationError("Topic cannot be empty")
    
    # Check stop before starting
    if stop_event and stop_event.is_set():
        raise GenerationStopped("Generation stopped by user")
    
    try:
        # Load configuration
        config = load_config()
        
        # Estimate slide count
        max_slides = estimate_slide_count_from_keywords(topic)
        max_slides = max(config["min_slides"], min(config["max_slides"], max_slides))
        
        # Build prompt
        prompt = config["adaptive_prompt_template"].format(
            system=config["system_prompt"],
            user_request=topic.strip(),
            min_points=config["min_points"],
            max_points=config["max_points"]
        )
        
        logger.info(f"Generating {max_slides} slides for topic: {topic[:50]}...")
        
        # Generate content with Ollama
        ollama_client = OllamaClient(timeout=config["ollama_timeout"])
        response = ollama_client.ask_ollama(prompt, config["model"], stop_event)
        
        # Check stop after LLM response
        if stop_event and stop_event.is_set():
            raise GenerationStopped("Generation stopped by user")
        
        # Parse JSON response
        slides_data = extract_and_parse_json(response)
        
        if not slides_data:
            # Check if this was due to stopping
            if stop_event and stop_event.is_set():
                raise GenerationStopped("Generation stopped by user")
                
            # Fallback: Create basic slides
            logger.warning("Failed to parse LLM response, creating fallback slides")
            slides_data = [
                {"title": topic.strip() or "Presentation", "points": []},
                {"title": "Overview", "points": [
                    "• Key concepts and fundamentals",
                    "• Practical applications and examples", 
                    "• Implementation strategies",
                    "• Best practices and considerations"
                ]}
            ]
        
        # Validate and coerce slides
        slides = validate_and_coerce_slides(slides_data, max_slides)
        
        # Check stop before creating presentation
        if stop_event and stop_event.is_set():
            raise GenerationStopped("Generation stopped by user")
        
        # Create presentation
        file_path = create_presentation(slides, template_path, font_name, font_size)
        
        # Metadata for client - SIMPLIFIED
        metadata = {
            "slide_count": len(slides),
            "topic": topic,
            "template_used": bool(template_path and os.path.exists(template_path))
        }
        
        logger.info(f"Successfully generated presentation with {len(slides)} slides")
        return file_path, metadata
        
    except GenerationStopped:
        raise
    except Exception as e:
        logger.error(f"Generation failed: {str(e)}")
        raise PPTGenerationError(f"Failed to generate presentation: {str(e)}")

# --------------------------
# Utility Functions
# --------------------------

def cleanup_temp_file(file_path: str) -> None:
    """Safely clean up temporary files"""
    try:
        if file_path and os.path.exists(file_path):
            os.unlink(file_path)
            logger.info(f"Cleaned up temporary file: {file_path}")
    except Exception as e:
        logger.warning(f"Failed to cleanup file {file_path}: {e}")

def validate_template(template_path: str) -> bool:
    """Validate if template file is a valid PowerPoint file"""
    if not os.path.exists(template_path):
        return False
    
    try:
        # Try to open as presentation
        prs = Presentation(template_path)
        return len(prs.slide_layouts) > 0
    except Exception:
        return False